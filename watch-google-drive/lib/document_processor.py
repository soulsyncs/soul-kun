"""
ドキュメント処理モジュール

各種ファイル形式からテキストを抽出し、チャンクに分割する機能を提供します。

使用例:
    from lib.document_processor import extract_text, TextChunker

    # テキスト抽出
    text = extract_text(file_content, 'pdf')

    # チャンク分割
    chunker = TextChunker(chunk_size=1000, chunk_overlap=200)
    chunks = chunker.split(text)

設計ドキュメント:
    docs/06_phase3_google_drive_integration.md
"""

import io
import re
import hashlib
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Optional
import logging

logger = logging.getLogger(__name__)


# ================================================================
# データクラス定義
# ================================================================

@dataclass
class ExtractedDocument:
    """抽出されたドキュメント情報"""
    text: str
    metadata: dict = field(default_factory=dict)
    pages: list[dict] = field(default_factory=list)
    total_pages: Optional[int] = None
    title: Optional[str] = None
    author: Optional[str] = None

    @property
    def char_count(self) -> int:
        return len(self.text)

    @property
    def text_hash(self) -> str:
        """テキストのSHA-256ハッシュ"""
        return hashlib.sha256(self.text.encode('utf-8')).hexdigest()


@dataclass
class Chunk:
    """チャンクデータ"""
    index: int                          # チャンク番号（0始まり）
    content: str                        # チャンクのテキスト
    char_count: int                     # 文字数
    start_position: int                 # 元文書での開始位置
    end_position: int                   # 元文書での終了位置
    page_number: Optional[int] = None   # ページ番号（PDFの場合）
    section_title: Optional[str] = None # セクションタイトル
    section_hierarchy: list[str] = field(default_factory=list)  # セクション階層

    @property
    def content_hash(self) -> str:
        """チャンクのSHA-256ハッシュ"""
        return hashlib.sha256(self.content.encode('utf-8')).hexdigest()


# ================================================================
# テキスト抽出クラス
# ================================================================

class TextExtractor(ABC):
    """テキスト抽出の基底クラス"""

    @abstractmethod
    def extract(self, content: bytes) -> str:
        """ファイルからテキストを抽出"""
        pass

    @abstractmethod
    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """テキストとメタデータを抽出"""
        pass


class PDFExtractor(TextExtractor):
    """PDFからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """PDFからテキストを抽出"""
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """PDFからテキストとメタデータを抽出"""
        import fitz

        doc = fitz.open(stream=content, filetype="pdf")

        pages = []
        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text()
            pages.append({
                "page_number": page_num,
                "text": page_text,
                "char_count": len(page_text)
            })

        metadata = {
            "total_pages": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "creator": doc.metadata.get("creator", ""),
            "creation_date": doc.metadata.get("creationDate", ""),
        }

        doc.close()

        return ExtractedDocument(
            text="\n".join([p["text"] for p in pages]),
            pages=pages,
            metadata=metadata,
            total_pages=len(pages),
            title=metadata.get("title"),
            author=metadata.get("author"),
        )


class DocxExtractor(TextExtractor):
    """Word文書（.docx）からテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """DOCXからテキストを抽出"""
        from docx import Document

        doc = Document(io.BytesIO(content))
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """DOCXからテキストとメタデータを抽出"""
        from docx import Document

        doc = Document(io.BytesIO(content))

        paragraphs = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append({
                    "index": i,
                    "text": para.text,
                    "style": para.style.name if para.style else None
                })

        # 見出しを抽出
        headings = [
            p for p in paragraphs
            if p["style"] and "Heading" in p["style"]
        ]

        metadata = {
            "total_paragraphs": len(paragraphs),
            "headings": headings,
            "author": doc.core_properties.author,
            "title": doc.core_properties.title,
            "created": str(doc.core_properties.created) if doc.core_properties.created else None,
            "modified": str(doc.core_properties.modified) if doc.core_properties.modified else None,
        }

        return ExtractedDocument(
            text="\n".join([p["text"] for p in paragraphs]),
            metadata=metadata,
            title=doc.core_properties.title,
            author=doc.core_properties.author,
        )


class DocExtractor(TextExtractor):
    """
    Word文書（.doc - バイナリ形式）からテキストを抽出

    .doc形式はMicrosoft Word 97-2003で使用されていたバイナリ形式です。
    antiwordコマンドを使用してテキストを抽出します。

    フォールバック戦略:
    1. antiwordでテキスト抽出を試行（軽量・高速）
    2. 失敗した場合、LibreOfficeで.docx変換を試行
    3. 変換されたファイルからpython-docxでテキスト抽出
    """

    # ファイルサイズの上限（50MB）
    MAX_FILE_SIZE = 50 * 1024 * 1024

    def extract(self, content: bytes) -> str:
        """DOCからテキストを抽出"""
        text, method = self._extract_with_fallback(content)
        if text is None:
            raise ValueError("DOCファイルからのテキスト抽出に失敗しました")
        logger.info(f"DOCテキスト抽出成功（方法: {method}）: {len(text)}文字")
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """DOCからテキストとメタデータを抽出"""
        text, method = self._extract_with_fallback(content)
        if text is None:
            raise ValueError("DOCファイルからのテキスト抽出に失敗しました")

        # 段落を抽出
        paragraphs = []
        for i, line in enumerate(text.split("\n")):
            if line.strip():
                paragraphs.append({
                    "index": i,
                    "text": line,
                    "style": None  # .doc形式ではスタイル情報を取得できない
                })

        metadata = {
            "total_paragraphs": len(paragraphs),
            "extraction_method": method,
            "format": "doc",
            "char_count": len(text)
        }

        logger.info(f"DOCテキスト抽出成功（方法: {method}）: {len(text)}文字")

        return ExtractedDocument(
            text=text,
            metadata=metadata,
        )

    def _extract_with_fallback(self, content: bytes) -> tuple[Optional[str], str]:
        """
        フォールバック戦略でテキストを抽出

        Returns:
            (抽出されたテキスト, 使用した方法)
        """
        import subprocess
        import tempfile
        import os

        # ファイルサイズチェック
        if len(content) > self.MAX_FILE_SIZE:
            logger.error(f"ファイルサイズが大きすぎます: {len(content) / 1024 / 1024:.2f}MB（上限: 50MB）")
            return (None, "failed")

        # 一時ファイルに書き込み
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            # 戦略1: antiwordでテキスト抽出
            logger.info("antiwordでテキスト抽出を試行")
            text = self._extract_with_antiword(tmp_path)
            if text:
                return (text, "antiword")

            # 戦略2: LibreOfficeで.docx変換
            logger.info("antiwordが失敗したため、LibreOfficeで変換を試行")
            text = self._extract_with_libreoffice(tmp_path)
            if text:
                return (text, "libreoffice")

            logger.error("全ての抽出方法が失敗しました")
            return (None, "failed")

        finally:
            # 一時ファイルを削除
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

    def _extract_with_antiword(self, file_path: str) -> Optional[str]:
        """antiwordを使用してテキストを抽出"""
        import subprocess

        try:
            result = subprocess.run(
                ["antiword", "-m", "UTF-8", file_path],
                capture_output=True,
                text=True,
                timeout=30,
                check=True
            )

            text = result.stdout
            if not text.strip():
                logger.warning("antiwordでテキストが抽出できませんでした")
                return None

            return text

        except subprocess.CalledProcessError as e:
            logger.warning(f"antiwordの実行に失敗: {e}")
            logger.debug(f"stderr: {e.stderr}")
            return None

        except FileNotFoundError:
            logger.warning("antiwordがインストールされていません")
            return None

        except subprocess.TimeoutExpired:
            logger.warning("antiwordがタイムアウトしました")
            return None

    def _extract_with_libreoffice(self, file_path: str) -> Optional[str]:
        """LibreOfficeで.docx変換してテキストを抽出"""
        import subprocess
        import tempfile
        import os
        from docx import Document

        try:
            # 一時出力ディレクトリ
            output_dir = tempfile.mkdtemp()

            try:
                # LibreOfficeで変換
                result = subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "docx",
                        "--outdir", output_dir,
                        file_path
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                    check=True
                )

                # 変換されたファイルのパスを構築
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                docx_path = os.path.join(output_dir, f"{base_name}.docx")

                if not os.path.exists(docx_path):
                    logger.warning(f"変換されたファイルが見つかりません: {docx_path}")
                    return None

                # python-docxでテキスト抽出
                doc = Document(docx_path)
                text = "\n".join([para.text for para in doc.paragraphs])

                return text if text.strip() else None

            finally:
                # 一時ディレクトリをクリーンアップ
                import shutil
                try:
                    shutil.rmtree(output_dir)
                except Exception:
                    pass

        except subprocess.CalledProcessError as e:
            logger.warning(f"LibreOfficeの実行に失敗: {e}")
            return None

        except FileNotFoundError:
            logger.warning("LibreOfficeがインストールされていません")
            return None

        except subprocess.TimeoutExpired:
            logger.warning("LibreOfficeがタイムアウトしました")
            return None

        except Exception as e:
            logger.warning(f"LibreOffice変換中にエラー: {e}")
            return None


class TextFileExtractor(TextExtractor):
    """テキストファイルからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """TXT/MDからテキストを抽出"""
        # UTF-8を試し、失敗したらShift-JISを試す
        for encoding in ["utf-8", "shift-jis", "cp932", "euc-jp"]:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        # 最後の手段: errors='replace'
        return content.decode("utf-8", errors="replace")

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """TXT/MDからテキストとメタデータを抽出"""
        text = self.extract(content)

        lines = text.split("\n")

        # Markdownの見出しを抽出
        headings = []
        for i, line in enumerate(lines):
            if line.startswith("#"):
                level = len(line) - len(line.lstrip("#"))
                heading_text = line.lstrip("#").strip()
                headings.append({
                    "line_number": i + 1,
                    "level": level,
                    "text": heading_text
                })

        return ExtractedDocument(
            text=text,
            metadata={
                "total_lines": len(lines),
                "total_chars": len(text),
                "headings": headings
            }
        )


class HTMLExtractor(TextExtractor):
    """HTMLからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """HTMLからテキストを抽出"""
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")

        # スクリプトとスタイルを除去
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        return soup.get_text(separator="\n", strip=True)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """HTMLからテキストとメタデータを抽出"""
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")

        # タイトルを取得
        title = soup.title.string if soup.title else ""

        # 見出しを抽出
        headings = []
        for level in range(1, 7):
            for heading in soup.find_all(f"h{level}"):
                headings.append({
                    "level": level,
                    "text": heading.get_text(strip=True)
                })

        # スクリプトとスタイルを除去
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)

        return ExtractedDocument(
            text=text,
            metadata={
                "title": title,
                "headings": headings,
                "total_chars": len(text)
            },
            title=title,
        )


class ExcelExtractor(TextExtractor):
    """Excelからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """XLSXからテキストを抽出"""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                row_texts = []
                for cell in row:
                    if cell.value is not None:
                        row_texts.append(str(cell.value))
                if row_texts:
                    texts.append("\t".join(row_texts))

        wb.close()
        return "\n".join(texts)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """XLSXからテキストとメタデータを抽出"""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        sheets = []
        all_text = []

        for sheet in wb.worksheets:
            sheet_text = []
            for row in sheet.iter_rows():
                row_texts = []
                for cell in row:
                    if cell.value is not None:
                        row_texts.append(str(cell.value))
                if row_texts:
                    sheet_text.append("\t".join(row_texts))

            sheets.append({
                "name": sheet.title,
                "text": "\n".join(sheet_text)
            })
            all_text.extend(sheet_text)

        wb.close()

        return ExtractedDocument(
            text="\n".join(all_text),
            metadata={
                "total_sheets": len(sheets),
                "sheets": sheets
            }
        )


class PowerPointExtractor(TextExtractor):
    """PowerPointからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """PPTXからテキストを抽出"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

        return "\n".join(texts)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """PPTXからテキストとメタデータを抽出"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        slides = []
        all_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides.append({
                "slide_number": slide_num,
                "text": "\n".join(slide_text)
            })
            all_text.extend(slide_text)

        return ExtractedDocument(
            text="\n".join(all_text),
            metadata={
                "total_slides": len(slides),
                "slides": slides
            }
        )


# ================================================================
# ファクトリ関数
# ================================================================

EXTRACTOR_MAP: dict[str, type[TextExtractor]] = {
    "pdf": PDFExtractor,
    "docx": DocxExtractor,
    "doc": DocExtractor,  # .doc形式専用のエクストラクター（antiword使用）
    "txt": TextFileExtractor,
    "md": TextFileExtractor,
    "html": HTMLExtractor,
    "htm": HTMLExtractor,
    "xlsx": ExcelExtractor,
    "xls": ExcelExtractor,
    "pptx": PowerPointExtractor,
    "ppt": PowerPointExtractor,
}


def get_extractor(file_type: str) -> Optional[TextExtractor]:
    """
    ファイル形式に対応するエクストラクターを取得

    Args:
        file_type: ファイル拡張子（小文字）

    Returns:
        TextExtractor または None
    """
    extractor_class = EXTRACTOR_MAP.get(file_type.lower())
    if extractor_class:
        return extractor_class()
    return None


def extract_text(content: bytes, file_type: str) -> str:
    """
    ファイルからテキストを抽出

    Args:
        content: ファイルのバイナリデータ
        file_type: ファイル拡張子

    Returns:
        抽出されたテキスト

    Raises:
        ValueError: サポートされていないファイル形式
    """
    extractor = get_extractor(file_type)
    if extractor is None:
        raise ValueError(f"サポートされていないファイル形式: {file_type}")
    return extractor.extract(content)


def extract_with_metadata(content: bytes, file_type: str) -> ExtractedDocument:
    """
    ファイルからテキストとメタデータを抽出

    Args:
        content: ファイルのバイナリデータ
        file_type: ファイル拡張子

    Returns:
        ExtractedDocument オブジェクト

    Raises:
        ValueError: サポートされていないファイル形式
    """
    extractor = get_extractor(file_type)
    if extractor is None:
        raise ValueError(f"サポートされていないファイル形式: {file_type}")
    return extractor.extract_with_metadata(content)


# ================================================================
# チャンク分割クラス
# ================================================================

class TextChunker:
    """
    テキストをチャンクに分割するクラス

    チャンク分割戦略:
    1. セマンティックな区切り（見出し、段落）を優先
    2. 文の途中で切らない
    3. オーバーラップで文脈を保持

    使用例:
        chunker = TextChunker(chunk_size=1000, chunk_overlap=200)
        chunks = chunker.split(text)
    """

    # セマンティックな区切り文字（優先度順）
    SEPARATORS = [
        "\n## ",      # Markdown H2
        "\n### ",     # Markdown H3
        "\n#### ",    # Markdown H4
        "\n\n",       # 空行（段落区切り）
        "\n",         # 改行
        "。",         # 日本語文末
        "．",         # 日本語ピリオド（全角）
        ". ",         # 英語文末
        "！",         # 日本語感嘆符
        "？",         # 日本語疑問符
        "! ",         # 英語感嘆符
        "? ",         # 英語疑問符
        "、",         # 日本語読点
        ", ",         # 英語コンマ
        " ",          # スペース
        "",           # 最後の手段（文字単位）
    ]

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100
    ):
        """
        Args:
            chunk_size: チャンクの最大文字数
            chunk_overlap: チャンク間のオーバーラップ文字数
            min_chunk_size: チャンクの最小文字数（これより短いチャンクは前のチャンクに結合）
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size

    def split(self, text: str) -> list[Chunk]:
        """
        テキストをチャンクに分割

        Args:
            text: 分割するテキスト

        Returns:
            チャンクのリスト
        """
        if not text or len(text) == 0:
            return []

        # 見出しを抽出
        headings = self._extract_headings(text)

        # 分割
        raw_chunks = self._split_text(text)

        # チャンクオブジェクトを作成
        chunks = []
        current_position = 0
        current_heading = None
        current_hierarchy: list[str] = []

        for i, chunk_text in enumerate(raw_chunks):
            # このチャンクの開始位置を検索
            chunk_start = text.find(chunk_text, current_position)
            if chunk_start == -1:
                chunk_start = current_position
            chunk_end = chunk_start + len(chunk_text)

            # このチャンクに含まれる見出しを更新
            for heading in headings:
                if heading["position"] >= chunk_start and heading["position"] < chunk_end:
                    current_heading = heading["text"]
                    current_hierarchy = heading["hierarchy"]

            chunk = Chunk(
                index=i,
                content=chunk_text,
                char_count=len(chunk_text),
                start_position=chunk_start,
                end_position=chunk_end,
                section_title=current_heading,
                section_hierarchy=current_hierarchy.copy() if current_hierarchy else []
            )
            chunks.append(chunk)

            current_position = max(chunk_end - self.chunk_overlap, chunk_start + 1)

        return chunks

    def split_with_pages(self, pages: list[dict]) -> list[Chunk]:
        """
        ページ情報付きのテキストをチャンクに分割（PDF用）

        Args:
            pages: [{"page_number": 1, "text": "..."}]

        Returns:
            チャンクのリスト（page_number付き）
        """
        chunks = []
        chunk_index = 0

        for page in pages:
            page_number = page["page_number"]
            page_text = page["text"]

            if not page_text.strip():
                continue

            # ページ内でチャンク分割
            page_chunks = self._split_text(page_text)

            current_position = 0
            for chunk_text in page_chunks:
                chunk_start = page_text.find(chunk_text, current_position)
                if chunk_start == -1:
                    chunk_start = current_position
                chunk_end = chunk_start + len(chunk_text)

                chunk = Chunk(
                    index=chunk_index,
                    content=chunk_text,
                    char_count=len(chunk_text),
                    start_position=chunk_start,
                    end_position=chunk_end,
                    page_number=page_number
                )
                chunks.append(chunk)

                chunk_index += 1
                current_position = max(chunk_end - self.chunk_overlap, chunk_start + 1)

        return chunks

    def _split_text(self, text: str) -> list[str]:
        """テキストを分割（再帰的）"""
        if len(text) <= self.chunk_size:
            return [text] if text.strip() else []

        # 各セパレータで分割を試みる
        for separator in self.SEPARATORS:
            if separator == "":
                # 最後の手段: 文字数で強制分割
                return self._split_by_length(text)

            if separator in text:
                splits = self._split_by_separator(text, separator)
                if len(splits) > 1:
                    # 再帰的に分割
                    result = []
                    for split in splits:
                        result.extend(self._split_text(split))
                    return self._merge_small_chunks(result)

        # どのセパレータでも分割できない場合
        return self._split_by_length(text)

    def _split_by_separator(self, text: str, separator: str) -> list[str]:
        """セパレータで分割"""
        splits = text.split(separator)

        # セパレータを復元（最後以外）
        result = []
        for i, split in enumerate(splits):
            if i < len(splits) - 1:
                result.append(split + separator)
            else:
                result.append(split)

        return [s for s in result if s.strip()]

    def _split_by_length(self, text: str) -> list[str]:
        """文字数で強制分割"""
        chunks = []
        step = self.chunk_size - self.chunk_overlap
        if step <= 0:
            step = self.chunk_size

        for i in range(0, len(text), step):
            chunk = text[i:i + self.chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        return chunks

    def _merge_small_chunks(self, chunks: list[str]) -> list[str]:
        """小さすぎるチャンクを前のチャンクに結合"""
        if not chunks:
            return chunks

        result = [chunks[0]]

        for chunk in chunks[1:]:
            if len(chunk) < self.min_chunk_size and result:
                # 前のチャンクと結合
                combined = result[-1] + chunk
                if len(combined) <= self.chunk_size:
                    result[-1] = combined
                else:
                    result.append(chunk)
            else:
                result.append(chunk)

        return result

    def _extract_headings(self, text: str) -> list[dict]:
        """見出しを抽出"""
        headings = []

        # Markdown形式の見出し
        md_heading_pattern = r'^(#{1,6})\s+(.+)$'
        for match in re.finditer(md_heading_pattern, text, re.MULTILINE):
            level = len(match.group(1))
            heading_text = match.group(2)
            position = match.start()

            # 階層を構築（簡易版）
            hierarchy = [heading_text]

            headings.append({
                "level": level,
                "text": heading_text,
                "position": position,
                "hierarchy": hierarchy
            })

        return headings


# ================================================================
# ドキュメントプロセッサ（統合クラス）
# ================================================================

class DocumentProcessor:
    """
    ドキュメント処理の統合クラス

    テキスト抽出からチャンク分割までを一括で処理します。

    使用例:
        processor = DocumentProcessor(chunk_size=1000, chunk_overlap=200)
        result = processor.process(file_content, 'pdf')
    """

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100,
    ):
        self.chunker = TextChunker(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            min_chunk_size=min_chunk_size,
        )

    def process(
        self,
        content: bytes,
        file_type: str,
    ) -> tuple[ExtractedDocument, list[Chunk]]:
        """
        ドキュメントを処理

        Args:
            content: ファイルのバイナリデータ
            file_type: ファイル拡張子

        Returns:
            (ExtractedDocument, チャンクのリスト)

        Raises:
            ValueError: サポートされていないファイル形式
        """
        # テキスト抽出
        doc = extract_with_metadata(content, file_type)

        # チャンク分割
        if doc.pages and file_type == 'pdf':
            # PDFの場合はページ情報を保持
            chunks = self.chunker.split_with_pages(doc.pages)
        else:
            chunks = self.chunker.split(doc.text)

        return doc, chunks

    def compute_file_hash(self, content: bytes) -> str:
        """ファイルのSHA-256ハッシュを計算"""
        return hashlib.sha256(content).hexdigest()


# ================================================================
# エクスポート
# ================================================================

__all__ = [
    # データクラス
    'ExtractedDocument',
    'Chunk',
    # テキスト抽出
    'TextExtractor',
    'PDFExtractor',
    'DocxExtractor',
    'DocExtractor',  # .doc形式専用
    'TextFileExtractor',
    'HTMLExtractor',
    'ExcelExtractor',
    'PowerPointExtractor',
    'get_extractor',
    'extract_text',
    'extract_with_metadata',
    'EXTRACTOR_MAP',
    # チャンク分割
    'TextChunker',
    # 統合クラス
    'DocumentProcessor',
]
