"""
ドキュメント処理モジュール

各種ファイル形式からテキストを抽出し、チャンクに分割する機能を提供します。

使用例:
    from lib.document_processor import extract_text, TextChunker

    # テキスト抽出
    text = extract_text(file_content, 'pdf')

    # チャンク分割
    chunker = TextChunker(chunk_size=1000, chunk_overlap=200)
    chunks = chunker.split(text)

設計ドキュメント:
    docs/06_phase3_google_drive_integration.md
"""

import io
import re
import hashlib
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Optional
import logging

logger = logging.getLogger(__name__)


# ================================================================
# チャンク品質フィルタリング（v10.13.2）
# ================================================================

# 目次を示すキーワード
TABLE_OF_CONTENTS_KEYWORDS = [
    "目　　次", "目　次", "目次",
    "もくじ", "CONTENTS", "Contents", "contents",
    "INDEX", "Index", "index",
    "TABLE OF CONTENTS", "Table of Contents",
]

# 目次ページの特徴的なパターン
TABLE_OF_CONTENTS_PATTERNS = [
    r"^第[一二三四五六七八九十\d]+章.*\d+$",  # 第一章　○○ 1
    r"^第[一二三四五六七八九十\d]+条.*\d+$",  # 第一条　○○ 1
    r"^[\d]+[\.\s]+.+[\s\.]+\d+$",              # 1. 概要 ... 3
    r"^[一二三四五六七八九十]+[\.\s、]+.+\d+$",  # 一、はじめに 5
    r"^【.+】.*\d+$",                           # 【総則】 1
    r"^（.+）.*\d+$",                           # （総則） 1
    r"^\d+[-\.\s].+[-\.\s]+\d+$",               # 1-1 概要...3
]

# ヘッダー/フッターの特徴的なパターン
HEADER_FOOTER_PATTERNS = [
    r"^[\d]+\s*/\s*[\d]+$",           # ページ番号: 1 / 10
    r"^-\s*\d+\s*-$",                  # ページ番号: - 1 -
    r"^page\s*\d+",                    # Page 1
    r"^\d+\s*ページ$",                 # 1ページ
    r"^©.*\d{4}",                      # © 2024
    r"^株式会社.*$",                   # 株式会社○○
    r"^confidential$",                 # CONFIDENTIAL
    r"^社外秘$",                       # 社外秘
    r"^取扱注意$",                     # 取扱注意
]

# 低品質コンテンツのパターン
LOW_QUALITY_PATTERNS = [
    r"^[\s\-_=\.]+$",                  # 区切り線のみ
    r"^[\d\.\s]+$",                     # 数字とスペースのみ
    r"^\s*$",                           # 空白のみ
]


def is_table_of_contents(text: str) -> bool:
    """
    目次ページかどうかを判定

    Args:
        text: チャンクのテキスト

    Returns:
        目次ページの場合True
    """
    if not text or len(text.strip()) == 0:
        return False

    # 1. 目次キーワードを含むか
    for keyword in TABLE_OF_CONTENTS_KEYWORDS:
        if keyword in text:
            # 目次キーワードを含む場合、目次パターンも多く含むか確認
            pattern_match_count = 0
            lines = text.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                for pattern in TABLE_OF_CONTENTS_PATTERNS:
                    if re.match(pattern, line):
                        pattern_match_count += 1
                        break

            # 全行の20%以上が目次パターンにマッチすれば目次と判定
            non_empty_lines = [l for l in lines if l.strip()]
            if non_empty_lines and pattern_match_count / len(non_empty_lines) > 0.2:
                return True

            # 目次キーワードが冒頭にある場合は目次と判定
            first_lines = "\n".join(lines[:5])
            if keyword in first_lines:
                return True

    # 2. 大部分が目次パターンにマッチするか（キーワードなしでも）
    lines = text.split("\n")
    pattern_match_count = 0
    for line in lines:
        line = line.strip()
        if not line:
            continue
        for pattern in TABLE_OF_CONTENTS_PATTERNS:
            if re.match(pattern, line):
                pattern_match_count += 1
                break

    non_empty_lines = [l for l in lines if l.strip()]
    if non_empty_lines and len(non_empty_lines) >= 5:
        # 50%以上が目次パターンにマッチすれば目次と判定
        if pattern_match_count / len(non_empty_lines) > 0.5:
            return True

    return False


def is_header_or_footer(text: str) -> bool:
    """
    ヘッダー/フッターかどうかを判定

    Args:
        text: テキスト

    Returns:
        ヘッダー/フッターの場合True
    """
    if not text:
        return False

    text = text.strip()

    # 短すぎるテキストはヘッダー/フッターの可能性が高い
    if len(text) < 10:
        for pattern in HEADER_FOOTER_PATTERNS:
            if re.match(pattern, text, re.IGNORECASE):
                return True

    return False


def is_low_quality_content(text: str) -> bool:
    """
    低品質コンテンツかどうかを判定

    Args:
        text: テキスト

    Returns:
        低品質コンテンツの場合True
    """
    if not text:
        return True

    text = text.strip()

    # 空または短すぎる
    if len(text) < 20:
        return True

    # 低品質パターンにマッチ
    for pattern in LOW_QUALITY_PATTERNS:
        if re.match(pattern, text):
            return True

    # 実質的な文字が少ない（記号や数字ばかり）
    japanese_chars = len(re.findall(r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]', text))
    alphabetic_chars = len(re.findall(r'[a-zA-Z]', text))
    meaningful_chars = japanese_chars + alphabetic_chars

    if len(text) > 0 and meaningful_chars / len(text) < 0.3:
        return True

    return False


def should_exclude_chunk(chunk_text: str) -> tuple[bool, str]:
    """
    チャンクを除外すべきかどうかを判定

    Args:
        chunk_text: チャンクのテキスト

    Returns:
        (除外すべきか, 理由)
    """
    # 1. 低品質コンテンツ
    if is_low_quality_content(chunk_text):
        return True, "low_quality_content"

    # 2. 目次ページ
    if is_table_of_contents(chunk_text):
        return True, "table_of_contents"

    # 3. 全体がヘッダー/フッターのみ
    lines = chunk_text.strip().split("\n")
    header_footer_lines = 0
    for line in lines:
        if is_header_or_footer(line):
            header_footer_lines += 1

    if lines and header_footer_lines / len(lines) > 0.8:
        return True, "header_footer_only"

    return False, ""


def calculate_chunk_quality_score(chunk_text: str) -> float:
    """
    チャンクの品質スコアを計算（0.0-1.0）

    スコアの基準:
    - 1.0: 高品質（本文、具体的な情報を含む）
    - 0.7-0.9: 中品質（見出し付きセクション）
    - 0.4-0.6: 低品質（断片的な情報）
    - 0.0-0.3: 除外対象（目次、ヘッダー等）

    Args:
        chunk_text: チャンクのテキスト

    Returns:
        品質スコア（0.0-1.0）
    """
    if not chunk_text:
        return 0.0

    score = 0.5  # ベーススコア

    # 1. 文字数による調整
    char_count = len(chunk_text.strip())
    if char_count < 50:
        score -= 0.3
    elif char_count < 100:
        score -= 0.1
    elif char_count > 500:
        score += 0.1

    # 2. 日本語/英語の意味のある文字の割合
    japanese_chars = len(re.findall(r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF]', chunk_text))
    alphabetic_chars = len(re.findall(r'[a-zA-Z]', chunk_text))
    meaningful_ratio = (japanese_chars + alphabetic_chars) / len(chunk_text) if chunk_text else 0

    if meaningful_ratio > 0.7:
        score += 0.2
    elif meaningful_ratio > 0.5:
        score += 0.1
    elif meaningful_ratio < 0.3:
        score -= 0.2

    # 3. 文の完結性（句点の存在）
    sentence_endings = len(re.findall(r'[。．.!?！？]', chunk_text))
    if sentence_endings >= 2:
        score += 0.1
    elif sentence_endings == 0:
        score -= 0.1

    # 4. 具体的な数値情報の存在
    has_specific_numbers = bool(re.search(r'\d+[日月年時分秒個件万円%]', chunk_text))
    if has_specific_numbers:
        score += 0.15

    # 5. 目次パターンのペナルティ
    toc_pattern_count = 0
    lines = chunk_text.split("\n")
    for line in lines:
        for pattern in TABLE_OF_CONTENTS_PATTERNS:
            if re.match(pattern, line.strip()):
                toc_pattern_count += 1
                break

    non_empty_lines = [l for l in lines if l.strip()]
    if non_empty_lines and toc_pattern_count / len(non_empty_lines) > 0.3:
        score -= 0.3

    # 6. 目次キーワードのペナルティ
    for keyword in TABLE_OF_CONTENTS_KEYWORDS:
        if keyword in chunk_text:
            score -= 0.2
            break

    # スコアを0.0-1.0の範囲にクランプ
    return max(0.0, min(1.0, score))


def extract_chunk_metadata(chunk_text: str) -> dict:
    """
    チャンクからメタデータを抽出

    Args:
        chunk_text: チャンクのテキスト

    Returns:
        メタデータ辞書
    """
    metadata = {}

    # 1. 条項番号の抽出
    article_match = re.search(r'第([一二三四五六七八九十\d]+)[条章節]', chunk_text)
    if article_match:
        metadata['article_number'] = article_match.group(0)

    # 2. カテゴリキーワードの抽出
    category_keywords = {
        '有給休暇': ['有給', '年休', '年次有給休暇', '休暇'],
        '給与': ['給与', '賃金', '給料', '報酬', '手当'],
        '勤務時間': ['勤務時間', '労働時間', '就業時間', '所定労働時間'],
        '休日': ['休日', '祝日', '週休'],
        '退職': ['退職', '退社', '解雇'],
        '懲戒': ['懲戒', '処分', '戒告'],
        '服務': ['服務', '規律'],
        '採用': ['採用', '入社', '雇用'],
        '福利厚生': ['福利厚生', '社会保険', '健康保険'],
        '経費': ['経費', '精算', '交通費'],
    }

    detected_categories = []
    for category, keywords in category_keywords.items():
        for keyword in keywords:
            if keyword in chunk_text:
                detected_categories.append(category)
                break

    if detected_categories:
        metadata['categories'] = detected_categories

    # 3. 数値情報の抽出
    # 日数
    days_match = re.search(r'(\d+)\s*日', chunk_text)
    if days_match:
        metadata['days_mentioned'] = int(days_match.group(1))

    # 金額
    amount_match = re.search(r'(\d+(?:,\d{3})*)\s*円', chunk_text)
    if amount_match:
        metadata['amount_mentioned'] = amount_match.group(1)

    # 4. 重要度のヒント
    importance_keywords = ['ただし', '注意', '重要', '禁止', '必ず', '厳禁']
    for keyword in importance_keywords:
        if keyword in chunk_text:
            metadata['has_important_note'] = True
            break

    return metadata


# ================================================================
# データクラス定義
# ================================================================

@dataclass
class ExtractedDocument:
    """抽出されたドキュメント情報"""
    text: str
    metadata: dict = field(default_factory=dict)
    pages: list[dict] = field(default_factory=list)
    total_pages: Optional[int] = None
    title: Optional[str] = None
    author: Optional[str] = None

    @property
    def char_count(self) -> int:
        return len(self.text)

    @property
    def text_hash(self) -> str:
        """テキストのSHA-256ハッシュ"""
        return hashlib.sha256(self.text.encode('utf-8')).hexdigest()


@dataclass
class Chunk:
    """チャンクデータ"""
    index: int                          # チャンク番号（0始まり）
    content: str                        # チャンクのテキスト
    char_count: int                     # 文字数
    start_position: int                 # 元文書での開始位置
    end_position: int                   # 元文書での終了位置
    page_number: Optional[int] = None   # ページ番号（PDFの場合）
    section_title: Optional[str] = None # セクションタイトル
    section_hierarchy: list[str] = field(default_factory=list)  # セクション階層
    quality_score: float = 1.0          # 品質スコア（0.0-1.0）（v10.13.2）
    chunk_metadata: dict = field(default_factory=dict)  # 追加メタデータ（v10.13.2）
    excluded: bool = False              # 除外フラグ（v10.13.2）
    exclusion_reason: str = ""          # 除外理由（v10.13.2）

    @property
    def content_hash(self) -> str:
        """チャンクのSHA-256ハッシュ"""
        return hashlib.sha256(self.content.encode('utf-8')).hexdigest()

    @property
    def is_high_quality(self) -> bool:
        """高品質チャンクかどうか（v10.13.2）"""
        return self.quality_score >= 0.5 and not self.excluded


# ================================================================
# テキスト抽出クラス
# ================================================================

class TextExtractor(ABC):
    """テキスト抽出の基底クラス"""

    @abstractmethod
    def extract(self, content: bytes) -> str:
        """ファイルからテキストを抽出"""
        pass

    @abstractmethod
    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """テキストとメタデータを抽出"""
        pass


class PDFExtractor(TextExtractor):
    """PDFからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """PDFからテキストを抽出"""
        import fitz  # PyMuPDF

        doc = fitz.open(stream=content, filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """PDFからテキストとメタデータを抽出"""
        import fitz

        doc = fitz.open(stream=content, filetype="pdf")

        pages = []
        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text()
            pages.append({
                "page_number": page_num,
                "text": page_text,
                "char_count": len(page_text)
            })

        metadata = {
            "total_pages": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
            "subject": doc.metadata.get("subject", ""),
            "creator": doc.metadata.get("creator", ""),
            "creation_date": doc.metadata.get("creationDate", ""),
        }

        doc.close()

        return ExtractedDocument(
            text="\n".join([p["text"] for p in pages]),
            pages=pages,
            metadata=metadata,
            total_pages=len(pages),
            title=metadata.get("title"),
            author=metadata.get("author"),
        )


class DocxExtractor(TextExtractor):
    """Word文書（.docx）からテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """DOCXからテキストを抽出"""
        from docx import Document

        doc = Document(io.BytesIO(content))
        text = "\n".join([para.text for para in doc.paragraphs])
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """DOCXからテキストとメタデータを抽出"""
        from docx import Document

        doc = Document(io.BytesIO(content))

        paragraphs = []
        for i, para in enumerate(doc.paragraphs):
            if para.text.strip():
                paragraphs.append({
                    "index": i,
                    "text": para.text,
                    "style": para.style.name if para.style else None
                })

        # 見出しを抽出
        headings = [
            p for p in paragraphs
            if p["style"] and "Heading" in p["style"]
        ]

        metadata = {
            "total_paragraphs": len(paragraphs),
            "headings": headings,
            "author": doc.core_properties.author,
            "title": doc.core_properties.title,
            "created": str(doc.core_properties.created) if doc.core_properties.created else None,
            "modified": str(doc.core_properties.modified) if doc.core_properties.modified else None,
        }

        return ExtractedDocument(
            text="\n".join([p["text"] for p in paragraphs]),
            metadata=metadata,
            title=doc.core_properties.title,
            author=doc.core_properties.author,
        )


class DocExtractor(TextExtractor):
    """
    Word文書（.doc - バイナリ形式）からテキストを抽出

    .doc形式はMicrosoft Word 97-2003で使用されていたバイナリ形式です。
    antiwordコマンドを使用してテキストを抽出します。

    フォールバック戦略:
    1. antiwordでテキスト抽出を試行（軽量・高速）
    2. 失敗した場合、LibreOfficeで.docx変換を試行
    3. 変換されたファイルからpython-docxでテキスト抽出
    """

    # ファイルサイズの上限（50MB）
    MAX_FILE_SIZE = 50 * 1024 * 1024

    def extract(self, content: bytes) -> str:
        """DOCからテキストを抽出"""
        text, method = self._extract_with_fallback(content)
        if text is None:
            raise ValueError("DOCファイルからのテキスト抽出に失敗しました")
        logger.info(f"DOCテキスト抽出成功（方法: {method}）: {len(text)}文字")
        return text

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """DOCからテキストとメタデータを抽出"""
        text, method = self._extract_with_fallback(content)
        if text is None:
            raise ValueError("DOCファイルからのテキスト抽出に失敗しました")

        # 段落を抽出
        paragraphs = []
        for i, line in enumerate(text.split("\n")):
            if line.strip():
                paragraphs.append({
                    "index": i,
                    "text": line,
                    "style": None  # .doc形式ではスタイル情報を取得できない
                })

        metadata = {
            "total_paragraphs": len(paragraphs),
            "extraction_method": method,
            "format": "doc",
            "char_count": len(text)
        }

        logger.info(f"DOCテキスト抽出成功（方法: {method}）: {len(text)}文字")

        return ExtractedDocument(
            text=text,
            metadata=metadata,
        )

    def _extract_with_fallback(self, content: bytes) -> tuple[Optional[str], str]:
        """
        フォールバック戦略でテキストを抽出

        Returns:
            (抽出されたテキスト, 使用した方法)
        """
        import subprocess
        import tempfile
        import os

        # ファイルサイズチェック
        if len(content) > self.MAX_FILE_SIZE:
            logger.error(f"ファイルサイズが大きすぎます: {len(content) / 1024 / 1024:.2f}MB（上限: 50MB）")
            return (None, "failed")

        # 一時ファイルに書き込み
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        try:
            # 戦略1: antiwordでテキスト抽出
            logger.info("antiwordでテキスト抽出を試行")
            text = self._extract_with_antiword(tmp_path)
            if text:
                return (text, "antiword")

            # 戦略2: LibreOfficeで.docx変換
            logger.info("antiwordが失敗したため、LibreOfficeで変換を試行")
            text = self._extract_with_libreoffice(tmp_path)
            if text:
                return (text, "libreoffice")

            logger.error("全ての抽出方法が失敗しました")
            return (None, "failed")

        finally:
            # 一時ファイルを削除
            try:
                os.unlink(tmp_path)
            except Exception:
                pass

    def _extract_with_antiword(self, file_path: str) -> Optional[str]:
        """antiwordを使用してテキストを抽出"""
        import subprocess

        try:
            result = subprocess.run(
                ["antiword", "-m", "UTF-8", file_path],
                capture_output=True,
                text=True,
                timeout=30,
                check=True
            )

            text = result.stdout
            if not text.strip():
                logger.warning("antiwordでテキストが抽出できませんでした")
                return None

            return text

        except subprocess.CalledProcessError as e:
            logger.warning(f"antiwordの実行に失敗: {e}")
            logger.debug(f"stderr: {e.stderr}")
            return None

        except FileNotFoundError:
            logger.warning("antiwordがインストールされていません")
            return None

        except subprocess.TimeoutExpired:
            logger.warning("antiwordがタイムアウトしました")
            return None

    def _extract_with_libreoffice(self, file_path: str) -> Optional[str]:
        """LibreOfficeで.docx変換してテキストを抽出"""
        import subprocess
        import tempfile
        import os
        from docx import Document

        try:
            # 一時出力ディレクトリ
            output_dir = tempfile.mkdtemp()

            try:
                # LibreOfficeで変換
                result = subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "docx",
                        "--outdir", output_dir,
                        file_path
                    ],
                    capture_output=True,
                    text=True,
                    timeout=60,
                    check=True
                )

                # 変換されたファイルのパスを構築
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                docx_path = os.path.join(output_dir, f"{base_name}.docx")

                if not os.path.exists(docx_path):
                    logger.warning(f"変換されたファイルが見つかりません: {docx_path}")
                    return None

                # python-docxでテキスト抽出
                doc = Document(docx_path)
                text = "\n".join([para.text for para in doc.paragraphs])

                return text if text.strip() else None

            finally:
                # 一時ディレクトリをクリーンアップ
                import shutil
                try:
                    shutil.rmtree(output_dir)
                except Exception:
                    pass

        except subprocess.CalledProcessError as e:
            logger.warning(f"LibreOfficeの実行に失敗: {e}")
            return None

        except FileNotFoundError:
            logger.warning("LibreOfficeがインストールされていません")
            return None

        except subprocess.TimeoutExpired:
            logger.warning("LibreOfficeがタイムアウトしました")
            return None

        except Exception as e:
            logger.warning(f"LibreOffice変換中にエラー: {e}")
            return None


class TextFileExtractor(TextExtractor):
    """テキストファイルからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """TXT/MDからテキストを抽出"""
        # UTF-8を試し、失敗したらShift-JISを試す
        for encoding in ["utf-8", "shift-jis", "cp932", "euc-jp"]:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        # 最後の手段: errors='replace'
        return content.decode("utf-8", errors="replace")

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """TXT/MDからテキストとメタデータを抽出"""
        text = self.extract(content)

        lines = text.split("\n")

        # Markdownの見出しを抽出
        headings = []
        for i, line in enumerate(lines):
            if line.startswith("#"):
                level = len(line) - len(line.lstrip("#"))
                heading_text = line.lstrip("#").strip()
                headings.append({
                    "line_number": i + 1,
                    "level": level,
                    "text": heading_text
                })

        return ExtractedDocument(
            text=text,
            metadata={
                "total_lines": len(lines),
                "total_chars": len(text),
                "headings": headings
            }
        )


class HTMLExtractor(TextExtractor):
    """HTMLからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """HTMLからテキストを抽出"""
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")

        # スクリプトとスタイルを除去
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        return soup.get_text(separator="\n", strip=True)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """HTMLからテキストとメタデータを抽出"""
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(content, "html.parser")

        # タイトルを取得
        title = soup.title.string if soup.title else ""

        # 見出しを抽出
        headings = []
        for level in range(1, 7):
            for heading in soup.find_all(f"h{level}"):
                headings.append({
                    "level": level,
                    "text": heading.get_text(strip=True)
                })

        # スクリプトとスタイルを除去
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)

        return ExtractedDocument(
            text=text,
            metadata={
                "title": title,
                "headings": headings,
                "total_chars": len(text)
            },
            title=title,
        )


class ExcelExtractor(TextExtractor):
    """Excelからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """XLSXからテキストを抽出"""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                row_texts = []
                for cell in row:
                    if cell.value is not None:
                        row_texts.append(str(cell.value))
                if row_texts:
                    texts.append("\t".join(row_texts))

        wb.close()
        return "\n".join(texts)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """XLSXからテキストとメタデータを抽出"""
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

        sheets = []
        all_text = []

        for sheet in wb.worksheets:
            sheet_text = []
            for row in sheet.iter_rows():
                row_texts = []
                for cell in row:
                    if cell.value is not None:
                        row_texts.append(str(cell.value))
                if row_texts:
                    sheet_text.append("\t".join(row_texts))

            sheets.append({
                "name": sheet.title,
                "text": "\n".join(sheet_text)
            })
            all_text.extend(sheet_text)

        wb.close()

        return ExtractedDocument(
            text="\n".join(all_text),
            metadata={
                "total_sheets": len(sheets),
                "sheets": sheets
            }
        )


class PowerPointExtractor(TextExtractor):
    """PowerPointからテキストを抽出"""

    def extract(self, content: bytes) -> str:
        """PPTXからテキストを抽出"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)

        return "\n".join(texts)

    def extract_with_metadata(self, content: bytes) -> ExtractedDocument:
        """PPTXからテキストとメタデータを抽出"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))

        slides = []
        all_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides.append({
                "slide_number": slide_num,
                "text": "\n".join(slide_text)
            })
            all_text.extend(slide_text)

        return ExtractedDocument(
            text="\n".join(all_text),
            metadata={
                "total_slides": len(slides),
                "slides": slides
            }
        )


# ================================================================
# ファクトリ関数
# ================================================================

EXTRACTOR_MAP: dict[str, type[TextExtractor]] = {
    "pdf": PDFExtractor,
    "docx": DocxExtractor,
    "doc": DocExtractor,  # .doc形式専用のエクストラクター（antiword使用）
    "txt": TextFileExtractor,
    "md": TextFileExtractor,
    "html": HTMLExtractor,
    "htm": HTMLExtractor,
    "xlsx": ExcelExtractor,
    "xls": ExcelExtractor,
    "pptx": PowerPointExtractor,
    "ppt": PowerPointExtractor,
}


def get_extractor(file_type: str) -> Optional[TextExtractor]:
    """
    ファイル形式に対応するエクストラクターを取得

    Args:
        file_type: ファイル拡張子（小文字）

    Returns:
        TextExtractor または None
    """
    extractor_class = EXTRACTOR_MAP.get(file_type.lower())
    if extractor_class:
        return extractor_class()
    return None


def extract_text(content: bytes, file_type: str) -> str:
    """
    ファイルからテキストを抽出

    Args:
        content: ファイルのバイナリデータ
        file_type: ファイル拡張子

    Returns:
        抽出されたテキスト

    Raises:
        ValueError: サポートされていないファイル形式
    """
    extractor = get_extractor(file_type)
    if extractor is None:
        raise ValueError(f"サポートされていないファイル形式: {file_type}")
    return extractor.extract(content)


def extract_with_metadata(content: bytes, file_type: str) -> ExtractedDocument:
    """
    ファイルからテキストとメタデータを抽出

    Args:
        content: ファイルのバイナリデータ
        file_type: ファイル拡張子

    Returns:
        ExtractedDocument オブジェクト

    Raises:
        ValueError: サポートされていないファイル形式
    """
    extractor = get_extractor(file_type)
    if extractor is None:
        raise ValueError(f"サポートされていないファイル形式: {file_type}")
    return extractor.extract_with_metadata(content)


# ================================================================
# チャンク分割クラス
# ================================================================

class TextChunker:
    """
    テキストをチャンクに分割するクラス

    チャンク分割戦略:
    1. セマンティックな区切り（見出し、段落）を優先
    2. 文の途中で切らない
    3. オーバーラップで文脈を保持

    使用例:
        chunker = TextChunker(chunk_size=1000, chunk_overlap=200)
        chunks = chunker.split(text)
    """

    # セマンティックな区切り文字（優先度順）
    SEPARATORS = [
        "\n## ",      # Markdown H2
        "\n### ",     # Markdown H3
        "\n#### ",    # Markdown H4
        "\n\n",       # 空行（段落区切り）
        "\n",         # 改行
        "。",         # 日本語文末
        "．",         # 日本語ピリオド（全角）
        ". ",         # 英語文末
        "！",         # 日本語感嘆符
        "？",         # 日本語疑問符
        "! ",         # 英語感嘆符
        "? ",         # 英語疑問符
        "、",         # 日本語読点
        ", ",         # 英語コンマ
        " ",          # スペース
        "",           # 最後の手段（文字単位）
    ]

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100
    ):
        """
        Args:
            chunk_size: チャンクの最大文字数
            chunk_overlap: チャンク間のオーバーラップ文字数
            min_chunk_size: チャンクの最小文字数（これより短いチャンクは前のチャンクに結合）
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size

    def split(self, text: str) -> list[Chunk]:
        """
        テキストをチャンクに分割

        Args:
            text: 分割するテキスト

        Returns:
            チャンクのリスト
        """
        if not text or len(text) == 0:
            return []

        # 見出しを抽出
        headings = self._extract_headings(text)

        # 分割
        raw_chunks = self._split_text(text)

        # チャンクオブジェクトを作成
        chunks = []
        current_position = 0
        current_heading = None
        current_hierarchy: list[str] = []

        for i, chunk_text in enumerate(raw_chunks):
            # このチャンクの開始位置を検索
            chunk_start = text.find(chunk_text, current_position)
            if chunk_start == -1:
                chunk_start = current_position
            chunk_end = chunk_start + len(chunk_text)

            # このチャンクに含まれる見出しを更新
            for heading in headings:
                if heading["position"] >= chunk_start and heading["position"] < chunk_end:
                    current_heading = heading["text"]
                    current_hierarchy = heading["hierarchy"]

            # 品質評価（v10.13.2）
            excluded, exclusion_reason = should_exclude_chunk(chunk_text)
            quality_score = calculate_chunk_quality_score(chunk_text)
            chunk_metadata = extract_chunk_metadata(chunk_text)

            chunk = Chunk(
                index=i,
                content=chunk_text,
                char_count=len(chunk_text),
                start_position=chunk_start,
                end_position=chunk_end,
                section_title=current_heading,
                section_hierarchy=current_hierarchy.copy() if current_hierarchy else [],
                quality_score=quality_score,
                chunk_metadata=chunk_metadata,
                excluded=excluded,
                exclusion_reason=exclusion_reason,
            )
            chunks.append(chunk)

            current_position = max(chunk_end - self.chunk_overlap, chunk_start + 1)

        return chunks

    def split_with_pages(self, pages: list[dict]) -> list[Chunk]:
        """
        ページ情報付きのテキストをチャンクに分割（PDF用）

        Args:
            pages: [{"page_number": 1, "text": "..."}]

        Returns:
            チャンクのリスト（page_number付き）
        """
        chunks = []
        chunk_index = 0

        for page in pages:
            page_number = page["page_number"]
            page_text = page["text"]

            if not page_text.strip():
                continue

            # ページ内でチャンク分割
            page_chunks = self._split_text(page_text)

            current_position = 0
            for chunk_text in page_chunks:
                chunk_start = page_text.find(chunk_text, current_position)
                if chunk_start == -1:
                    chunk_start = current_position
                chunk_end = chunk_start + len(chunk_text)

                # 品質評価（v10.13.2）
                excluded, exclusion_reason = should_exclude_chunk(chunk_text)
                quality_score = calculate_chunk_quality_score(chunk_text)
                chunk_metadata = extract_chunk_metadata(chunk_text)

                chunk = Chunk(
                    index=chunk_index,
                    content=chunk_text,
                    char_count=len(chunk_text),
                    start_position=chunk_start,
                    end_position=chunk_end,
                    page_number=page_number,
                    quality_score=quality_score,
                    chunk_metadata=chunk_metadata,
                    excluded=excluded,
                    exclusion_reason=exclusion_reason,
                )
                chunks.append(chunk)

                chunk_index += 1
                current_position = max(chunk_end - self.chunk_overlap, chunk_start + 1)

        return chunks

    def _split_text(self, text: str) -> list[str]:
        """テキストを分割（再帰的）"""
        if len(text) <= self.chunk_size:
            return [text] if text.strip() else []

        # 各セパレータで分割を試みる
        for separator in self.SEPARATORS:
            if separator == "":
                # 最後の手段: 文字数で強制分割
                return self._split_by_length(text)

            if separator in text:
                splits = self._split_by_separator(text, separator)
                if len(splits) > 1:
                    # 再帰的に分割
                    result = []
                    for split in splits:
                        result.extend(self._split_text(split))
                    return self._merge_small_chunks(result)

        # どのセパレータでも分割できない場合
        return self._split_by_length(text)

    def _split_by_separator(self, text: str, separator: str) -> list[str]:
        """セパレータで分割"""
        splits = text.split(separator)

        # セパレータを復元（最後以外）
        result = []
        for i, split in enumerate(splits):
            if i < len(splits) - 1:
                result.append(split + separator)
            else:
                result.append(split)

        return [s for s in result if s.strip()]

    def _split_by_length(self, text: str) -> list[str]:
        """文字数で強制分割"""
        chunks = []
        step = self.chunk_size - self.chunk_overlap
        if step <= 0:
            step = self.chunk_size

        for i in range(0, len(text), step):
            chunk = text[i:i + self.chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        return chunks

    def _merge_small_chunks(self, chunks: list[str]) -> list[str]:
        """小さすぎるチャンクを前のチャンクに結合"""
        if not chunks:
            return chunks

        result = [chunks[0]]

        for chunk in chunks[1:]:
            if len(chunk) < self.min_chunk_size and result:
                # 前のチャンクと結合
                combined = result[-1] + chunk
                if len(combined) <= self.chunk_size:
                    result[-1] = combined
                else:
                    result.append(chunk)
            else:
                result.append(chunk)

        return result

    def _extract_headings(self, text: str) -> list[dict]:
        """見出しを抽出"""
        headings = []

        # Markdown形式の見出し
        md_heading_pattern = r'^(#{1,6})\s+(.+)$'
        for match in re.finditer(md_heading_pattern, text, re.MULTILINE):
            level = len(match.group(1))
            heading_text = match.group(2)
            position = match.start()

            # 階層を構築（簡易版）
            hierarchy = [heading_text]

            headings.append({
                "level": level,
                "text": heading_text,
                "position": position,
                "hierarchy": hierarchy
            })

        return headings


# ================================================================
# ドキュメントプロセッサ（統合クラス）
# ================================================================

class DocumentProcessor:
    """
    ドキュメント処理の統合クラス

    テキスト抽出からチャンク分割までを一括で処理します。

    使用例:
        processor = DocumentProcessor(chunk_size=1000, chunk_overlap=200)
        result = processor.process(file_content, 'pdf')

    品質フィルタリング（v10.13.2）:
        processor = DocumentProcessor(
            chunk_size=1000,
            chunk_overlap=200,
            filter_low_quality=True,  # 低品質チャンクを除外
            min_quality_score=0.4      # 品質スコアの閾値
        )
    """

    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100,
        filter_low_quality: bool = True,   # 低品質チャンクを除外（v10.13.2）
        min_quality_score: float = 0.4,    # 品質スコアの閾値（v10.13.2）
    ):
        self.chunker = TextChunker(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            min_chunk_size=min_chunk_size,
        )
        self.filter_low_quality = filter_low_quality
        self.min_quality_score = min_quality_score

    def process(
        self,
        content: bytes,
        file_type: str,
    ) -> tuple[ExtractedDocument, list[Chunk]]:
        """
        ドキュメントを処理

        Args:
            content: ファイルのバイナリデータ
            file_type: ファイル拡張子

        Returns:
            (ExtractedDocument, チャンクのリスト)

        Raises:
            ValueError: サポートされていないファイル形式
        """
        # テキスト抽出
        doc = extract_with_metadata(content, file_type)

        # チャンク分割
        if doc.pages and file_type == 'pdf':
            # PDFの場合はページ情報を保持
            chunks = self.chunker.split_with_pages(doc.pages)
        else:
            chunks = self.chunker.split(doc.text)

        return doc, chunks

    def process_with_quality_filter(
        self,
        content: bytes,
        file_type: str,
    ) -> tuple[ExtractedDocument, list[Chunk], list[Chunk]]:
        """
        ドキュメントを処理し、品質フィルタリングを適用（v10.13.2）

        Args:
            content: ファイルのバイナリデータ
            file_type: ファイル拡張子

        Returns:
            (ExtractedDocument, 高品質チャンクのリスト, 除外されたチャンクのリスト)

        Raises:
            ValueError: サポートされていないファイル形式
        """
        doc, all_chunks = self.process(content, file_type)

        if not self.filter_low_quality:
            return doc, all_chunks, []

        # 品質フィルタリング
        high_quality_chunks = []
        excluded_chunks = []

        for chunk in all_chunks:
            if chunk.excluded:
                excluded_chunks.append(chunk)
                logger.debug(
                    f"チャンク除外 [index={chunk.index}]: {chunk.exclusion_reason}"
                )
            elif chunk.quality_score < self.min_quality_score:
                chunk.excluded = True
                chunk.exclusion_reason = f"low_quality_score_{chunk.quality_score:.2f}"
                excluded_chunks.append(chunk)
                logger.debug(
                    f"チャンク除外 [index={chunk.index}]: 品質スコア {chunk.quality_score:.2f} < {self.min_quality_score}"
                )
            else:
                high_quality_chunks.append(chunk)

        # インデックスを再割り当て
        for i, chunk in enumerate(high_quality_chunks):
            chunk.index = i

        logger.info(
            f"品質フィルタリング完了: "
            f"全{len(all_chunks)}チャンク → "
            f"高品質{len(high_quality_chunks)}チャンク "
            f"(除外{len(excluded_chunks)}チャンク)"
        )

        return doc, high_quality_chunks, excluded_chunks

    def get_quality_report(self, chunks: list[Chunk]) -> dict:
        """
        チャンクの品質レポートを生成（v10.13.2）

        Args:
            chunks: チャンクのリスト

        Returns:
            品質レポート辞書
        """
        if not chunks:
            return {
                "total_chunks": 0,
                "high_quality_chunks": 0,
                "excluded_chunks": 0,
                "average_quality_score": 0.0,
                "exclusion_reasons": {},
            }

        high_quality = [c for c in chunks if c.is_high_quality]
        excluded = [c for c in chunks if c.excluded]

        exclusion_reasons: dict[str, int] = {}
        for chunk in excluded:
            reason = chunk.exclusion_reason or "unknown"
            exclusion_reasons[reason] = exclusion_reasons.get(reason, 0) + 1

        return {
            "total_chunks": len(chunks),
            "high_quality_chunks": len(high_quality),
            "excluded_chunks": len(excluded),
            "average_quality_score": sum(c.quality_score for c in chunks) / len(chunks),
            "exclusion_reasons": exclusion_reasons,
            "quality_score_distribution": {
                "0.0-0.2": len([c for c in chunks if c.quality_score < 0.2]),
                "0.2-0.4": len([c for c in chunks if 0.2 <= c.quality_score < 0.4]),
                "0.4-0.6": len([c for c in chunks if 0.4 <= c.quality_score < 0.6]),
                "0.6-0.8": len([c for c in chunks if 0.6 <= c.quality_score < 0.8]),
                "0.8-1.0": len([c for c in chunks if c.quality_score >= 0.8]),
            },
        }

    def compute_file_hash(self, content: bytes) -> str:
        """ファイルのSHA-256ハッシュを計算"""
        return hashlib.sha256(content).hexdigest()


# ================================================================
# エクスポート
# ================================================================

__all__ = [
    # データクラス
    'ExtractedDocument',
    'Chunk',
    # テキスト抽出
    'TextExtractor',
    'PDFExtractor',
    'DocxExtractor',
    'DocExtractor',  # .doc形式専用
    'TextFileExtractor',
    'HTMLExtractor',
    'ExcelExtractor',
    'PowerPointExtractor',
    'get_extractor',
    'extract_text',
    'extract_with_metadata',
    'EXTRACTOR_MAP',
    # チャンク分割
    'TextChunker',
    # 統合クラス
    'DocumentProcessor',
    # 品質フィルタリング（v10.13.2）
    'is_table_of_contents',
    'is_header_or_footer',
    'is_low_quality_content',
    'should_exclude_chunk',
    'calculate_chunk_quality_score',
    'extract_chunk_metadata',
    'TABLE_OF_CONTENTS_KEYWORDS',
    'TABLE_OF_CONTENTS_PATTERNS',
]
